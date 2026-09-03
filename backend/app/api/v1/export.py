"""
Export endpoints for generating PPTX and PDF from slides.
Uses python-pptx for PowerPoint and ReportLab for PDF.
"""
import io
import httpx
from fastapi import APIRouter, HTTPException
from fastapi.responses import Response
from pydantic import BaseModel
from typing import List, Optional, Any

router = APIRouter()


class ExportRequest(BaseModel):
    presentation_id: Optional[str] = None
    title: Optional[str] = "Untitled Presentation"
    slides: List[dict] = []
    theme: Optional[dict] = None
    include_images: bool = True


async def download_image(url: str) -> Optional[bytes]:
    """Download image bytes for embedding into PPT/PDF."""
    if not url or not isinstance(url, str):
        return None
    if url.startswith("data:image"):
        # Base64 data URL
        try:
            import base64
            header, encoded = url.split(",", 1)
            return base64.b64decode(encoded)
        except Exception:
            return None
    try:
        async with httpx.AsyncClient(timeout=10.0, follow_redirects=True) as client:
            res = await client.get(url)
            if res.status_code == 200:
                return res.content
    except Exception as e:
        print(f"[Export] Image download failed: {e}")
    return None


@router.post("/pptx")
async def export_pptx(req: ExportRequest):
    """Export presentation to PowerPoint (.pptx) format."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dgm.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
    except ImportError:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt, Emu
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE
        except ImportError:
            raise HTTPException(status_code=500, detail="python-pptx not installed. Run: pip install python-pptx")

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        theme = req.theme or {}
        primary_hex = (theme.get("primary") or "#0077B6").lstrip("#")
        primary_rgb = RGBColor(int(primary_hex[0:2], 16), int(primary_hex[2:4], 16), int(primary_hex[4:6], 16))

        for idx, slide_data in enumerate(req.slides):
            blank_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(blank_layout)

            layout_type = slide_data.get("layout", "bullets")
            title_text = slide_data.get("title", "")

            # Background image
            image_url = slide_data.get("image_url")
            if req.include_images and image_url:
                img_bytes = await download_image(image_url)
                if img_bytes:
                    try:
                        img_stream = io.BytesIO(img_bytes)
                        slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    except Exception as e:
                        print(f"[Export] PPT image insert failed: {e}")

            # Semi-transparent dark overlay for text readability
            overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            overlay.fill.solid()
            overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
            overlay.fill.transparency = 0.5
            overlay.line.fill.background()

            # Title text box
            title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(12), Inches(1.5))
            title_frame = title_box.text_frame
            title_frame.word_wrap = True
            title_para = title_frame.paragraphs[0]
            title_run = title_para.add_run()
            title_run.text = title_text
            title_run.font.size = Pt(44)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(255, 255, 255)

            # Layout-specific content
            if layout_type == "title":
                subtitle = slide_data.get("subtitle", "")
                if subtitle:
                    sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.8), Inches(12), Inches(1))
                    sub_frame = sub_box.text_frame
                    sub_run = sub_frame.paragraphs[0].add_run()
                    sub_run.text = subtitle
                    sub_run.font.size = Pt(22)
                    sub_run.font.color.rgb = RGBColor(220, 220, 220)

            elif layout_type == "bullets":
                bullets = slide_data.get("bullets", [])
                bullet_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(12), Inches(4.5))
                bullet_frame = bullet_box.text_frame
                bullet_frame.word_wrap = True
                for i, b in enumerate(bullets[:6]):
                    p = bullet_frame.add_paragraph() if i > 0 else bullet_frame.paragraphs[0]
                    r = p.add_run()
                    r.text = f"• {b}"
                    r.font.size = Pt(18)
                    r.font.color.rgb = RGBColor(255, 255, 255)
                    p.space_after = Pt(10)

            elif layout_type == "metrics":
                metrics = slide_data.get("metrics", [])
                metric_colors = [RGBColor(34, 211, 238), RGBColor(192, 132, 252), RGBColor(52, 211, 153)]
                for i, m in enumerate(metrics[:3]):
                    x_pos = Inches(0.7 + i * 4.2)
                    m_box = slide.shapes.add_textbox(x_pos, Inches(3), Inches(4), Inches(2))
                    m_frame = m_box.text_frame
                    m_frame.word_wrap = True

                    val_p = m_frame.paragraphs[0]
                    val_r = val_p.add_run()
                    val_r.text = m.get("value", "")
                    val_r.font.size = Pt(52)
                    val_r.font.bold = True
                    val_r.font.color.rgb = metric_colors[i % len(metric_colors)]

                    lbl_p = m_frame.add_paragraph()
                    lbl_r = lbl_p.add_run()
                    lbl_r.text = m.get("label", "")
                    lbl_r.font.size = Pt(14)
                    lbl_r.font.color.rgb = RGBColor(220, 220, 220)

            elif layout_type == "two_column":
                left_col = slide_data.get("left_column", {})
                right_col = slide_data.get("right_column", {})

                # Left column
                left_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.5), Inches(6), Inches(4.5))
                left_frame = left_box.text_frame
                left_frame.word_wrap = True
                lt_run = left_frame.paragraphs[0].add_run()
                lt_run.text = left_col.get("title", "")
                lt_run.font.size = Pt(20)
                lt_run.font.bold = True
                lt_run.font.color.rgb = RGBColor(251, 113, 133)
                for pt in left_col.get("points", [])[:4]:
                    p = left_frame.add_paragraph()
                    r = p.add_run()
                    r.text = f"• {pt}"
                    r.font.size = Pt(14)
                    r.font.color.rgb = RGBColor(255, 255, 255)

                # Right column
                right_box = slide.shapes.add_textbox(Inches(6.9), Inches(2.5), Inches(6), Inches(4.5))
                right_frame = right_box.text_frame
                right_frame.word_wrap = True
                rt_run = right_frame.paragraphs[0].add_run()
                rt_run.text = right_col.get("title", "")
                rt_run.font.size = Pt(20)
                rt_run.font.bold = True
                rt_run.font.color.rgb = RGBColor(52, 211, 153)
                for pt in right_col.get("points", [])[:4]:
                    p = right_frame.add_paragraph()
                    r = p.add_run()
                    r.text = f"• {pt}"
                    r.font.size = Pt(14)
                    r.font.color.rgb = RGBColor(255, 255, 255)

            elif layout_type == "quote":
                quote = slide_data.get("quote", "")
                author = slide_data.get("quote_author", "")
                q_box = slide.shapes.add_textbox(Inches(1.5), Inches(3), Inches(10.3), Inches(3))
                q_frame = q_box.text_frame
                q_frame.word_wrap = True
                q_p = q_frame.paragraphs[0]
                q_p.alignment = 2  # PP_ALIGN.CENTER
                q_r = q_p.add_run()
                q_r.text = f'"{quote}"'
                q_r.font.size = Pt(28)
                q_r.font.italic = True
                q_r.font.color.rgb = RGBColor(255, 255, 255)

                a_p = q_frame.add_paragraph()
                a_p.alignment = 2
                a_r = a_p.add_run()
                a_r.text = f"— {author}"
                a_r.font.size = Pt(16)
                a_r.font.color.rgb = RGBColor(200, 200, 200)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        return Response(
            content=buf.read(),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f"attachment; filename={req.title}.pptx"}
        )

    except Exception as e:
        print(f"[Export] PPTX error: {e}")
        raise HTTPException(status_code=500, detail=f"PPTX export failed: {str(e)}")


@router.post("/pdf")
async def export_pdf(req: ExportRequest):
    """Export presentation to PDF format."""
    try:
        from reportlab.lib.pagesizes import landscape, letter
        from reportlab.pdfgen import canvas
        from reportlab.lib.colors import HexColor, black, white
        from reportlab.lib.utils import ImageReader
    except ImportError:
        raise HTTPException(status_code=500, detail="reportlab not installed. Run: pip install reportlab")

    try:
        buf = io.BytesIO()
        page_size = landscape(letter)
        c = canvas.Canvas(buf, pagesize=page_size)
        width, height = page_size

        theme = req.theme or {}
        primary_hex = theme.get("primary") or "#0077B6"

        for slide_data in req.slides:
            layout_type = slide_data.get("layout", "bullets")
            title = slide_data.get("title", "")

            # Background image
            image_url = slide_data.get("image_url")
            if req.include_images and image_url:
                img_bytes = await download_image(image_url)
                if img_bytes:
                    try:
                        img = ImageReader(io.BytesIO(img_bytes))
                        c.drawImage(img, 0, 0, width=width, height=height, preserveAspectRatio=False, mask="auto")
                    except Exception as e:
                        print(f"[Export] PDF image insert failed: {e}")

            # Dark overlay
            c.setFillColor(HexColor("#000000"))
            c.setFillAlpha(0.5)
            c.rect(0, 0, width, height, fill=True, stroke=False)
            c.setFillAlpha(1)

            # Title
            c.setFillColor(white)
            c.setFont("Helvetica-Bold", 32)
            c.drawString(50, height - 80, title[:80])

            # Layout content
            if layout_type == "title":
                subtitle = slide_data.get("subtitle", "")
                c.setFont("Helvetica", 18)
                c.setFillColor(HexColor("#DDDDDD"))
                c.drawString(50, height - 130, subtitle[:100])

            elif layout_type == "bullets":
                bullets = slide_data.get("bullets", [])
                c.setFont("Helvetica", 14)
                c.setFillColor(white)
                y = height - 180
                for b in bullets[:6]:
                    c.drawString(70, y, f"• {b[:100]}")
                    y -= 40

            elif layout_type == "metrics":
                metrics = slide_data.get("metrics", [])
                x_start = 80
                for i, m in enumerate(metrics[:3]):
                    x = x_start + i * (width / 3.5)
                    c.setFont("Helvetica-Bold", 48)
                    c.setFillColor(HexColor(primary_hex))
                    c.drawString(x, height / 2, m.get("value", ""))
                    c.setFont("Helvetica", 12)
                    c.setFillColor(HexColor("#DDDDDD"))
                    c.drawString(x, height / 2 - 30, m.get("label", ""))

            elif layout_type == "two_column":
                left = slide_data.get("left_column", {})
                right = slide_data.get("right_column", {})
                c.setFont("Helvetica-Bold", 16)
                c.setFillColor(HexColor("#FB7185"))
                c.drawString(60, height - 180, left.get("title", ""))
                c.setFont("Helvetica", 12)
                c.setFillColor(white)
                y = height - 220
                for pt in left.get("points", [])[:4]:
                    c.drawString(70, y, f"• {pt[:80]}")
                    y -= 30

                c.setFont("Helvetica-Bold", 16)
                c.setFillColor(HexColor("#34D399"))
                c.drawString(width / 2 + 30, height - 180, right.get("title", ""))
                c.setFont("Helvetica", 12)
                c.setFillColor(white)
                y = height - 220
                for pt in right.get("points", [])[:4]:
                    c.drawString(width / 2 + 40, y, f"• {pt[:80]}")
                    y -= 30

            elif layout_type == "quote":
                quote = slide_data.get("quote", "")
                author = slide_data.get("quote_author", "")
                c.setFont("Helvetica-Oblique", 22)
                c.setFillColor(white)
                # Simple center-ish quote
                c.drawCentredString(width / 2, height / 2 + 20, f'"{quote[:80]}"')
                c.setFont("Helvetica", 14)
                c.setFillColor(HexColor("#CCCCCC"))
                c.drawCentredString(width / 2, height / 2 - 30, f"— {author}")

            c.showPage()

        c.save()
        buf.seek(0)

        return Response(
            content=buf.read(),
            media_type="application/pdf",
            headers={"Content-Disposition": f"attachment; filename={req.title}.pdf"}
        )

    except Exception as e:
        print(f"[Export] PDF error: {e}")
        raise HTTPException(status_code=500, detail=f"PDF export failed: {str(e)}")
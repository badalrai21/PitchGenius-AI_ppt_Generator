import io
import logging
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

logger = logging.getLogger(__name__)

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def hex_to_rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color="#03045E", alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = hex_to_rgb(color)
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color="#F0F9FF"):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def build_pptx(
    title: str,
    slides_data: List[Dict[str, Any]],
    theme: Dict[str, str],
    images: Dict[int, bytes] = None
) -> bytes:
    """
    Build a complete .pptx file from structured slide data.
    Returns raw .pptx bytes ready for download.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    primary = theme.get("primary", "#0077B6")
    secondary = theme.get("secondary", "#00B4D8")
    accent = theme.get("accent", "#90E0EF")
    bg_color = theme.get("bg", "#FAFDFF")
    text_color = theme.get("text", "#03045E")
    muted_color = theme.get("muted", "#64748B")

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for idx, slide_data in enumerate(slides_data):
        slide = prs.slides.add_slide(blank_layout)
        layout = slide_data.get("layout", "bullets")

        # Background fill
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(bg_color)

        # Add image if available
        if images and idx in images and images[idx]:
            try:
                img_stream = io.BytesIO(images[idx])
                slide.shapes.add_picture(img_stream, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
                # Add dark overlay for readability
                overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT)
                overlay.fill.solid()
                overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
                overlay.fill.fore_color.brightness = 0.0
                overlay.line.fill.background()
                # Set transparency via XML
                from pptx.oxml.ns import qn
                solidFill = overlay.fill._fill.find(qn('a:solidFill'))
                if solidFill is not None:
                    srgb = solidFill.find(qn('a:srgbClr'))
                    if srgb is not None:
                        alpha = srgb.makeelement(qn('a:alpha'), {})
                        alpha.set('val', '40000')  # 40% opacity
                        srgb.append(alpha)
            except Exception as e:
                logger.warning(f"Could not embed image for slide {idx}: {e}")

        # ---- TITLE SLIDE ----
        if layout == "title":
            # Accent bar
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.08))
            bar.fill.solid()
            bar.fill.fore_color.rgb = hex_to_rgb(primary)
            bar.line.fill.background()

            add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.5),
                        slide_data.get("title", title), font_size=44, bold=True,
                        color=primary if bg_color == "#FAFDFF" or bg_color == "#FFFFFF" else "#FFFFFF",
                        alignment=PP_ALIGN.CENTER)

            if slide_data.get("subtitle"):
                add_textbox(slide, Inches(2), Inches(3.2), Inches(9), Inches(1),
                            slide_data["subtitle"], font_size=22,
                            color=muted_color, alignment=PP_ALIGN.CENTER)

            if slide_data.get("body"):
                add_textbox(slide, Inches(2.5), Inches(4.5), Inches(8), Inches(1),
                            slide_data["body"], font_size=14,
                            color=muted_color, alignment=PP_ALIGN.CENTER)

        # ---- METRICS SLIDE ----
        elif layout == "metrics":
            add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                        slide_data.get("title", ""), font_size=32, bold=True, color=text_color)

            metrics = slide_data.get("metrics", [])
            if metrics:
                col_width = Inches(3.5)
                gap = Inches(0.5)
                start_x = Inches(0.8)
                for m_idx, metric in enumerate(metrics[:3]):
                    x = start_x + (col_width + gap) * m_idx
                    card = add_rounded_rect(slide, x, Inches(2.5), col_width, Inches(3), fill_color="#F0F9FF")
                    add_textbox(slide, x, Inches(3), col_width, Inches(1.2),
                                metric.get("value", ""), font_size=48, bold=True,
                                color=primary, alignment=PP_ALIGN.CENTER)
                    add_textbox(slide, x, Inches(4.2), col_width, Inches(0.6),
                                metric.get("label", ""), font_size=14,
                                color=muted_color, alignment=PP_ALIGN.CENTER)

        # ---- TWO COLUMN SLIDE ----
        elif layout == "two_column":
            add_textbox(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                        slide_data.get("title", ""), font_size=32, bold=True, color=text_color)

            left_col = slide_data.get("left_column", {})
            right_col = slide_data.get("right_column", {})

            # Left card
            add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.5), fill_color="#F0F9FF")
            add_textbox(slide, Inches(1.2), Inches(2), Inches(4.8), Inches(0.6),
                        left_col.get("title", "Overview"), font_size=18, bold=True, color=primary)
            left_points = left_col.get("points", [])
            for p_idx, pt in enumerate(left_points[:5]):
                add_textbox(slide, Inches(1.2), Inches(2.8 + p_idx * 0.6), Inches(4.8), Inches(0.5),
                            f"✓  {pt}", font_size=13, color=text_color)

            # Right card
            add_rounded_rect(slide, Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), fill_color="#E0F7FA")
            add_textbox(slide, Inches(7.2), Inches(2), Inches(4.8), Inches(0.6),
                        right_col.get("title", "Solution"), font_size=18, bold=True, color=secondary)
            right_points = right_col.get("points", [])
            for p_idx, pt in enumerate(right_points[:5]):
                add_textbox(slide, Inches(7.2), Inches(2.8 + p_idx * 0.6), Inches(4.8), Inches(0.5),
                            f"→  {pt}", font_size=13, color=text_color)

        # ---- QUOTE SLIDE ----
        elif layout == "quote":
            quote_text = slide_data.get("quote", slide_data.get("body", ""))
            add_textbox(slide, Inches(1.5), Inches(2), Inches(10), Inches(2.5),
                        f'"{quote_text}"', font_size=28, bold=False,
                        color=text_color, alignment=PP_ALIGN.CENTER)
            if slide_data.get("quote_author"):
                add_textbox(slide, Inches(1.5), Inches(4.8), Inches(10), Inches(0.6),
                            f"— {slide_data['quote_author']}", font_size=14, bold=True,
                            color=primary, alignment=PP_ALIGN.CENTER)

        # ---- DEFAULT BULLETS SLIDE ----
        else:
            # Top accent line
            bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.5), Inches(1.5), Inches(0.06))
            bar.fill.solid()
            bar.fill.fore_color.rgb = hex_to_rgb(primary)
            bar.line.fill.background()

            add_textbox(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.8),
                        slide_data.get("title", ""), font_size=32, bold=True, color=text_color)

            if slide_data.get("subtitle"):
                add_textbox(slide, Inches(0.8), Inches(1.5), Inches(10), Inches(0.5),
                            slide_data["subtitle"], font_size=14, color=muted_color)

            bullets = slide_data.get("bullets", [])
            if not bullets and slide_data.get("body"):
                bullets = [b.strip() for b in slide_data["body"].split("\n") if b.strip()]

            for b_idx, bullet in enumerate(bullets[:6]):
                y = Inches(2.3 + b_idx * 0.75)
                card = add_rounded_rect(slide, Inches(0.8), y, Inches(11.5), Inches(0.6), fill_color="#F8FDFF")
                add_textbox(slide, Inches(1.2), y + Inches(0.08), Inches(10.8), Inches(0.45),
                            bullet, font_size=15, color=text_color)

        # Speaker notes
        if slide_data.get("speaker_notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["speaker_notes"]

    # Save to bytes
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read()